import os
import io
import json
import re
import uuid
import time
import sqlite3
import secrets
import asyncio
from datetime import date

import requests
from fastapi import FastAPI, UploadFile, File, Form, HTTPException, Body, Request
from fastapi.responses import FileResponse, JSONResponse, StreamingResponse
from fastapi.staticfiles import StaticFiles
from pypdf import PdfReader
from docx import Document
from docx.shared import Pt
from pptx import Presentation
from fpdf import FPDF
from dotenv import load_dotenv

load_dotenv()

# ============================================================
# КОНФИГ GIGACHAT
# ============================================================
GIGACHAT_AUTH_KEY = os.getenv("GIGACHAT_AUTH_KEY", "")
GIGACHAT_SCOPE = os.getenv("GIGACHAT_SCOPE", "GIGACHAT_API_PERS")

OAUTH_URL = "https://ngw.devices.sberbank.ru:9443/api/v2/oauth"
CHAT_URL = "https://gigachat.devices.sberbank.ru/api/v1/chat/completions"

# GigaChat использует сертификаты российского удостоверяющего центра.
# Для быстрого старта проверку SSL отключаем — см. README про продакшен-вариант.
VERIFY_SSL = False

# Максимальная длина текста лекции, которую отправляем в GigaChat за раз.
MAX_INPUT_CHARS = 15000

# Максимальный размер загружаемого файла (в байтах). 10 МБ с запасом
# хватает на любую текстовую лекцию — так пользователь не положит сервер,
# закинув видео или огромный архив по ошибке.
MAX_FILE_SIZE_BYTES = 10 * 1024 * 1024  # 10 МБ


class QuotaExceeded(Exception):
    """Внутреннее исключение: у модели закончился бесплатный баланс токенов."""
    pass


# Порядок моделей для разных задач — если на первой в списке кончился
# баланс, автоматически пробуем следующую. Карточки — самая частая
# операция, поэтому впереди Lite (у него больше всего токенов). Для
# уточнений и карты лекции первым ставим Pro — там важнее качество ответа.
MODELS_FOR_CARDS = ["GigaChat", "GigaChat-Pro"]
MODELS_FOR_CLARIFY = ["GigaChat-Pro", "GigaChat"]

# Тариф GigaChat (Freemium) разрешает только 1 одновременный запрос ко
# всему API-ключу. Если два пользователя нажмут кнопку одновременно —
# второй запрос без этой защиты просто получит ошибку от GigaChat.
# Семафор ставит все запросы в очередь: они выполняются по одному,
# автоматически дожидаясь освобождения "потока", вместо падения с ошибкой.
gigachat_semaphore = asyncio.Semaphore(1)

# Простой счётчик использований в памяти (сбрасывается при перезапуске
# сервера). Для реальной статистики между перезапусками потребуется база
# данных, но для теста на группе этого достаточно.
_usage_stats = {"date": str(date.today()), "count": 0}

_token_cache = {"access_token": None, "expires_at": 0}

# ============================================================
# ХРАНИЛИЩЕ ДЛЯ РАСШАРИВАЕМЫХ НАБОРОВ КАРТОЧЕК (SQLite)
# ============================================================
# ВАЖНО: на бесплатном тарифе Render диск не постоянный — файл базы
# может обнулиться при перезапуске/передеплое сервиса. Для теста на
# группе этого достаточно; для долгосрочного хранения ссылок в будущем
# потребуется подключить постоянную БД (например, Render Postgres).
DB_PATH = os.path.join(os.path.dirname(__file__), "cards.db")


def init_db():
    conn = sqlite3.connect(DB_PATH)
    conn.execute(
        """
        CREATE TABLE IF NOT EXISTS card_sets (
            id TEXT PRIMARY KEY,
            cards_json TEXT NOT NULL,
            title TEXT,
            created_at TEXT NOT NULL
        )
        """
    )
    conn.execute(
        """
        CREATE TABLE IF NOT EXISTS ip_usage (
            ip TEXT NOT NULL,
            date TEXT NOT NULL,
            count INTEGER NOT NULL DEFAULT 0,
            PRIMARY KEY (ip, date)
        )
        """
    )
    conn.commit()
    conn.close()


init_db()

# Бесплатный дневной лимит генераций на один IP-адрес. Пока нет аккаунтов —
# это самый простой честный способ ограничить бесплатное использование,
# не заставляя людей регистрироваться.
# Промокоды для безлимитного доступа (снимают дневной лимит + открывают
# платные фичи — кастомизацию и доп. форматы экспорта). Пока нет
# настоящей оплаты — коды выдаёшь вручную (себе и первым клиентам через
# Telegram). Задаются через .env / Environment Variables на Render,
# через запятую: PREMIUM_CODES=код1,код2,код3
PREMIUM_CODES = set(
    code.strip() for code in os.getenv("PREMIUM_CODES", "").split(",") if code.strip()
)


def is_premium(request: Request) -> bool:
    code = request.headers.get("x-premium-code", "")
    return code in PREMIUM_CODES


FREE_DAILY_LIMIT = 5


def get_client_ip(request: Request) -> str:
    """Render (и большинство хостингов) работают через прокси — реальный
    IP пользователя лежит в заголовке X-Forwarded-For, а не в
    request.client.host (там будет IP самого прокси-сервера)."""
    forwarded = request.headers.get("x-forwarded-for")
    if forwarded:
        return forwarded.split(",")[0].strip()
    return request.client.host if request.client else "unknown"


def check_and_use_quota(ip: str) -> int:
    """Проверяет дневной лимит для IP, увеличивает счётчик, если лимит не
    исчерпан. Возвращает, сколько генераций осталось на сегодня после
    этого запроса. Кидает 429, если лимит уже исчерпан."""
    today = str(date.today())
    conn = sqlite3.connect(DB_PATH)
    row = conn.execute(
        "SELECT count FROM ip_usage WHERE ip = ? AND date = ?", (ip, today)
    ).fetchone()
    used = row[0] if row else 0

    if used >= FREE_DAILY_LIMIT:
        conn.close()
        raise HTTPException(
            status_code=429,
            detail=(
                f"Бесплатный дневной лимит ({FREE_DAILY_LIMIT} генераций) исчерпан. "
                "Лимит обновится завтра, либо оформи безлимитный доступ."
            ),
        )

    conn.execute(
        "INSERT INTO ip_usage (ip, date, count) VALUES (?, ?, 1) "
        "ON CONFLICT(ip, date) DO UPDATE SET count = count + 1",
        (ip, today),
    )
    conn.commit()
    conn.close()
    return FREE_DAILY_LIMIT - (used + 1)


def get_remaining_quota(ip: str) -> int:
    """Только читает остаток лимита, не тратит его — для отображения в UI."""
    today = str(date.today())
    conn = sqlite3.connect(DB_PATH)
    row = conn.execute(
        "SELECT count FROM ip_usage WHERE ip = ? AND date = ?", (ip, today)
    ).fetchone()
    conn.close()
    used = row[0] if row else 0
    return max(0, FREE_DAILY_LIMIT - used)


def get_gigachat_token() -> str:
    if not GIGACHAT_AUTH_KEY:
        raise HTTPException(
            status_code=500,
            detail="GIGACHAT_AUTH_KEY не задан. Добавь его в .env файл.",
        )

    now = time.time()
    if _token_cache["access_token"] and _token_cache["expires_at"] > now + 5:
        return _token_cache["access_token"]

    headers = {
        "Content-Type": "application/x-www-form-urlencoded",
        "Accept": "application/json",
        "RqUID": str(uuid.uuid4()),
        "Authorization": f"Basic {GIGACHAT_AUTH_KEY}",
    }
    data = {"scope": GIGACHAT_SCOPE}

    try:
        resp = requests.post(OAUTH_URL, headers=headers, data=data, verify=VERIFY_SSL, timeout=30)
    except requests.exceptions.Timeout:
        raise HTTPException(status_code=504, detail="GigaChat не ответил за 30 секунд (получение токена). Попробуй ещё раз.")
    except requests.exceptions.ConnectionError:
        raise HTTPException(status_code=502, detail="Не удалось подключиться к серверу GigaChat. Проверь интернет-соединение.")

    if resp.status_code != 200:
        raise HTTPException(
            status_code=502,
            detail=f"Ошибка получения токена GigaChat: {resp.status_code} {resp.text}",
        )
    payload = resp.json()
    _token_cache["access_token"] = payload["access_token"]
    expires_at_raw = payload.get("expires_at", 0)
    _token_cache["expires_at"] = expires_at_raw / 1000 if expires_at_raw > 10**12 else now + 1700
    return _token_cache["access_token"]


def ask_gigachat(prompt: str, model: str = "GigaChat", max_tokens: int = 8000, temperature: float = 0.3):
    """Низкоуровневый запрос к конкретной модели GigaChat."""
    token = get_gigachat_token()
    headers = {
        "Content-Type": "application/json",
        "Accept": "application/json",
        "Authorization": f"Bearer {token}",
    }
    body = {
        "model": model,
        "messages": [
            {
                "role": "system",
                "content": (
                    "Ты помощник, который помогает студентам разбираться в лекциях. "
                    "Следуй формату, который просит пользователь, точно и без отступлений."
                ),
            },
            {"role": "user", "content": prompt},
        ],
        "temperature": temperature,
        "max_tokens": max_tokens,
    }

    try:
        resp = requests.post(CHAT_URL, headers=headers, json=body, verify=VERIFY_SSL, timeout=60)
    except requests.exceptions.Timeout:
        raise HTTPException(
            status_code=504,
            detail="GigaChat не ответил за 60 секунд. Возможно, лекция слишком большая — попробуй сократить текст.",
        )
    except requests.exceptions.ConnectionError:
        raise HTTPException(
            status_code=502,
            detail="Не удалось подключиться к серверу GigaChat. Проверь интернет-соединение и повтори попытку.",
        )

    # 402/429 обычно значит "закончился баланс токенов на эту модель"
    if resp.status_code in (402, 429):
        raise QuotaExceeded(f"Лимит токенов исчерпан для модели {model}: {resp.status_code} {resp.text}")

    if resp.status_code != 200:
        raise HTTPException(
            status_code=502,
            detail=f"Ошибка запроса к GigaChat: {resp.status_code} {resp.text}",
        )
    data = resp.json()
    choice = data["choices"][0]
    content = choice["message"]["content"]
    finish_reason = choice.get("finish_reason", "unknown")
    return content, finish_reason


def ask_gigachat_with_fallback(prompt: str, models: list[str], max_tokens: int = 8000, temperature: float = 0.3):
    """Пробует модели по очереди из списка — если на текущей закончился
    баланс токенов (QuotaExceeded), автоматически переходит к следующей.
    Так лимит одной модели не останавливает работу приложения."""
    last_error = None
    for model in models:
        try:
            return ask_gigachat(prompt, model=model, max_tokens=max_tokens, temperature=temperature)
        except QuotaExceeded as e:
            last_error = e
            continue
    raise HTTPException(
        status_code=402,
        detail=f"Лимит токенов исчерпан на всех доступных моделях ({', '.join(models)}). {last_error}",
    )


def extract_json(raw_text: str, finish_reason: str = "unknown"):
    """GigaChat иногда оборачивает ответ в ```json ... ``` — чистим и парсим.
    Модель иногда путает формат и вместо массива [] возвращает объект {}
    с ключами "0", "1", "2"... — приводим оба варианта к единому списку.
    Если ответ обрезан по лимиту токенов, пытаемся восстановить хотя бы те
    карточки, что успели прийти целиком, вместо полного отказа."""
    cleaned = raw_text.strip()
    cleaned = re.sub(r"^```(json)?", "", cleaned.strip())
    cleaned = re.sub(r"```$", "", cleaned.strip())
    cleaned = cleaned.strip()

    match = re.search(r"[\[{].*", cleaned, re.DOTALL)
    if match:
        cleaned = match.group(0)

    opening_char = cleaned[:1]
    closing_char = "]" if opening_char == "[" else "}"

    parsed = None
    try:
        parsed = json.loads(cleaned, strict=False)
    except json.JSONDecodeError:
        last_complete = cleaned.rfind("}")
        if last_complete == -1:
            raise HTTPException(
                status_code=502,
                detail=f"Модель вернула пустой или нечитаемый ответ. Сырой ответ: {raw_text[:500]}",
            )
        repaired = cleaned[: last_complete + 1] + closing_char
        try:
            parsed = json.loads(repaired, strict=False)
        except json.JSONDecodeError as e:
            if finish_reason == "length":
                reason_hint = "Не хватило max_tokens — ответ обрезан по лимиту длины."
            else:
                reason_hint = f"Модель завершила ответ некорректно (finish_reason: {finish_reason})."
            # Показываем именно то место, где сломался парсинг, а не всегда
            # начало строки — так виднее реальная причина (например,
            # неэкранированные кавычки внутри значения).
            pos = e.pos if hasattr(e, "pos") else 0
            window_start = max(0, pos - 150)
            window_end = min(len(raw_text), pos + 150)
            context = raw_text[window_start:window_end]
            raise HTTPException(
                status_code=502,
                detail=(
                    f"Ответ от модели повреждён и не удалось восстановить: {e}. {reason_hint} "
                    f"Фрагмент рядом с ошибкой: ...{context}..."
                ),
            )

    # Модель вернула объект вместо массива — берём значения по порядку ключей
    if isinstance(parsed, dict):
        parsed = list(parsed.values())

    return parsed


# ============================================================
# ИЗВЛЕЧЕНИЕ ТЕКСТА ИЗ ФАЙЛОВ
# ============================================================
def extract_text_from_upload(filename: str, content: bytes) -> str:
    ext = os.path.splitext(filename)[1].lower()
    text = ""
    stream = io.BytesIO(content)

    if ext in (".txt", ".md"):
        text = content.decode("utf-8", errors="ignore")
    elif ext == ".pdf":
        reader = PdfReader(stream)
        for page in reader.pages:
            extracted = page.extract_text()
            if extracted:
                text += extracted + "\n"
    elif ext == ".docx":
        doc = Document(stream)
        for para in doc.paragraphs:
            if para.text.strip():
                text += para.text + "\n"
    elif ext == ".pptx":
        prs = Presentation(stream)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        p_text = "".join(run.text for run in para.runs)
                        if p_text.strip():
                            text += p_text + "\n"
    else:
        raise HTTPException(status_code=400, detail=f"Формат {ext} не поддерживается")

    return text


# ============================================================
# ПРОМПТ ДЛЯ ГЕНЕРАЦИИ КАРТОЧЕК
# ============================================================
def build_prompt(source_text: str, level: str = "detailed") -> str:
    trimmed = source_text[:MAX_INPUT_CHARS]

    if level == "brief":
        answer_instruction = 'краткий ответ (1 предложение, только суть)'
    else:
        answer_instruction = 'подробный ответ (2-4 предложения с пояснением)'

    return (
        "Преврати текст лекции ниже в набор обучающих карточек для повторения. "
        "Каждая карточка — это один ключевой факт, определение или концепция.\n\n"
        "Верни ТОЛЬКО JSON-массив объектов со следующими полями:\n"
        '- "title": короткий заголовок темы карточки (3-6 слов)\n'
        '- "question": вопрос для самопроверки\n'
        f'- "answer": {answer_instruction}\n'
        '- "tags": массив из 1-3 ключевых слов темы\n\n'
        "КРИТИЧЕСКИ ВАЖНО: верни строго валидный JSON-МАССИВ (начинается с "
        "символа [ и заканчивается символом ]), а НЕ объект с ключами-номерами "
        '(НЕ { "0": ..., "1": ... }). Каждый элемент массива — это объект в '
        "фигурных скобках {}, а не строка в кавычках.\n\n"
        "ВАЖНО ПРО КАВЫЧКИ: если в тексте лекции встречаются названия или "
        "цитаты в кавычках (например, название проекта \"Pixel Peak\"), "
        "НЕ используй двойные кавычки внутри значений полей — заменяй их "
        "на одинарные кавычки ' или просто убирай, иначе сломаешь формат JSON.\n\n"
        "Сделай от 5 до 12 карточек в зависимости от объёма материала. "
        "Не добавляй ничего, кроме самого JSON-массива. Не форматируй JSON "
        "с отступами и переносами строк — верни компактную запись в одну строку, "
        "это экономит место в ответе.\n\n"
        f"Текст лекции:\n{trimmed}"
    )


# ============================================================
# FASTAPI ПРИЛОЖЕНИЕ
# ============================================================
app = FastAPI(title="Лекция -> Карточки")


def build_clarify_prompt(card_question: str, card_answer: str, user_question: str) -> str:
    """Промпт для уточняющего вопроса по конкретной карточке. Даём модели
    только контекст этой карточки (не всю лекцию) — так дешевле и быстрее,
    а модель отвечает по существу, не расплываясь."""
    return (
        "Студент изучает карточку для повторения со следующим содержимым:\n"
        f"Вопрос карточки: {card_question}\n"
        f"Ответ карточки: {card_answer}\n\n"
        f"Студент задал уточняющий вопрос: {user_question}\n\n"
        "Ответь кратко и по существу (2-4 предложения), не повторяя "
        "дословно то, что уже написано в ответе карточки — дай именно "
        "уточнение или объяснение сверх уже известного."
    )


@app.post("/api/generate-cards")
async def generate_cards(
    request: Request,
    file: UploadFile | None = File(default=None),
    text: str | None = Form(default=None),
    level: str = Form(default="detailed"),
):
    client_ip = get_client_ip(request)
    user_is_premium = is_premium(request)
    remaining_after = None if user_is_premium else check_and_use_quota(client_ip)

    if file is not None:
        content = await file.read()
        if len(content) > MAX_FILE_SIZE_BYTES:
            raise HTTPException(
                status_code=413,
                detail=f"Файл слишком большой (максимум {MAX_FILE_SIZE_BYTES // (1024*1024)} МБ)",
            )
        source_text = extract_text_from_upload(file.filename, content)
    elif text and text.strip():
        source_text = text
    else:
        raise HTTPException(status_code=400, detail="Нужно прислать файл или текст")

    if not source_text.strip():
        raise HTTPException(status_code=400, detail="Не удалось извлечь текст из файла")

    was_truncated = len(source_text) > MAX_INPUT_CHARS

    prompt = build_prompt(source_text, level=level)
    async with gigachat_semaphore:
        raw_response, finish_reason = ask_gigachat_with_fallback(prompt, models=MODELS_FOR_CARDS)
    cards = extract_json(raw_response, finish_reason)

    # Автоматически сохраняем каждую генерацию в БД — это и есть основа
    # истории: не нужно отдельного действия "сохранить", ссылка и запись
    # в истории появляются сразу при генерации.
    set_id = secrets.token_urlsafe(5)
    title = cards[0].get("title", "Карточки") if cards else "Карточки"
    conn = sqlite3.connect(DB_PATH)
    conn.execute(
        "INSERT INTO card_sets (id, cards_json, title, created_at) VALUES (?, ?, ?, ?)",
        (set_id, json.dumps(cards, ensure_ascii=False), title, str(date.today())),
    )
    conn.commit()
    conn.close()

    today = str(date.today())
    if _usage_stats["date"] != today:
        _usage_stats["date"] = today
        _usage_stats["count"] = 0
    _usage_stats["count"] += 1

    return JSONResponse({
        "cards": cards,
        "source_length": len(source_text),
        "max_input_chars": MAX_INPUT_CHARS,
        "was_truncated": was_truncated,
        "id": set_id,
        "url": f"/s/{set_id}",
        "title": title,
        "quota_remaining": remaining_after,
        "quota_limit": FREE_DAILY_LIMIT,
        "is_premium": user_is_premium,
    })


@app.get("/api/quota")
async def get_quota(request: Request):
    """Позволяет фронтенду показать остаток лимита до генерации, не тратя его."""
    if is_premium(request):
        return JSONResponse({"remaining": None, "limit": FREE_DAILY_LIMIT, "is_premium": True})
    client_ip = get_client_ip(request)
    remaining = get_remaining_quota(client_ip)
    return JSONResponse({"remaining": remaining, "limit": FREE_DAILY_LIMIT, "is_premium": False})


@app.get("/api/limits")
async def get_limits():
    return JSONResponse({
        "max_input_chars": MAX_INPUT_CHARS,
        "max_file_size_mb": MAX_FILE_SIZE_BYTES // (1024 * 1024),
    })


@app.get("/api/set/{set_id}")
async def get_set(set_id: str):
    conn = sqlite3.connect(DB_PATH)
    row = conn.execute(
        "SELECT cards_json, title FROM card_sets WHERE id = ?", (set_id,)
    ).fetchone()
    conn.close()

    if not row:
        raise HTTPException(status_code=404, detail="Набор карточек не найден (возможно, ссылка устарела)")

    cards_json, title = row
    return JSONResponse({"cards": json.loads(cards_json), "title": title})


@app.get("/s/{set_id}")
async def view_shared_set(set_id: str):
    return FileResponse(os.path.join("static", "shared.html"))


@app.get("/api/stats")
async def get_stats():
    """Счётчик генераций за сегодня — пригодится, когда будешь вводить тарифы."""
    today = str(date.today())
    if _usage_stats["date"] != today:
        return JSONResponse({"date": today, "count": 0})
    return JSONResponse(_usage_stats)


@app.post("/api/clarify")
async def clarify_card(payload: dict = Body(...)):
    """Отвечает на уточняющий вопрос студента по конкретной карточке."""
    card_question = payload.get("card_question", "")
    card_answer = payload.get("card_answer", "")
    user_question = payload.get("user_question", "")

    if not user_question.strip():
        raise HTTPException(status_code=400, detail="Вопрос пустой")

    prompt = build_clarify_prompt(card_question, card_answer, user_question)
    async with gigachat_semaphore:
        answer_text, _ = ask_gigachat_with_fallback(
            prompt, models=MODELS_FOR_CLARIFY, max_tokens=500
        )

    return JSONResponse({"answer": answer_text.strip()})


@app.get("/robots.txt")
async def robots_txt():
    return FileResponse(os.path.join("static", "robots.txt"), media_type="text/plain")


@app.get("/sitemap.xml")
async def sitemap_xml():
    return FileResponse(os.path.join("static", "sitemap.xml"), media_type="application/xml")


@app.post("/api/export/{fmt}")
async def export_cards(fmt: str, request: Request, payload: dict = Body(...)):
    """Экспорт набора карточек в docx или pdf — платная фича, доступна
    только с валидным промокодом (см. is_premium)."""
    if not is_premium(request):
        raise HTTPException(
            status_code=403,
            detail="Экспорт в этот формат доступен только с промокодом безлимита.",
        )

    cards = payload.get("cards", [])
    if not cards:
        raise HTTPException(status_code=400, detail="Нет карточек для экспорта")

    if fmt == "docx":
        buffer = build_docx(cards)
        media_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        ascii_name, utf8_name = "cards.docx", "карточки.docx"
    elif fmt == "pdf":
        buffer = build_pdf(cards)
        media_type = "application/pdf"
        ascii_name, utf8_name = "cards.pdf", "карточки.pdf"
    else:
        raise HTTPException(status_code=400, detail="Неизвестный формат (доступны: docx, pdf)")

    # HTTP-заголовки должны быть latin-1 — кириллицу напрямую в имя файла
    # положить нельзя. Даём ASCII-имя как основное (filename=) и правильно
    # закодированное кириллическое имя (filename*=UTF-8''...) для браузеров,
    # которые его поддерживают — так скачается "карточки.pdf", а не "cards.pdf".
    from urllib.parse import quote
    content_disposition = (
        f'attachment; filename="{ascii_name}"; '
        f"filename*=UTF-8''{quote(utf8_name)}"
    )

    return StreamingResponse(
        buffer,
        media_type=media_type,
        headers={"Content-Disposition": content_disposition},
    )


def build_docx(cards: list) -> io.BytesIO:
    doc = Document()
    title = doc.add_heading("Карточки для повторения", level=1)

    for card in cards:
        doc.add_heading(card.get("title", ""), level=2)
        p_tags = doc.add_paragraph()
        p_tags.add_run(" · ".join(card.get("tags", []))).italic = True

        p_q = doc.add_paragraph()
        run_q = p_q.add_run("Вопрос: " + card.get("question", ""))
        run_q.bold = True

        doc.add_paragraph("Ответ: " + card.get("answer", ""))
        doc.add_paragraph("")  # пустая строка между карточками

    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer


FONT_PATH = os.path.join(os.path.dirname(__file__), "fonts", "DejaVuSans.ttf")


def build_pdf(cards: list) -> io.BytesIO:
    pdf = FPDF()
    pdf.add_page()
    # Стандартные PDF-шрифты не поддерживают кириллицу — подключаем
    # DejaVu Sans, у которого есть нужные символы.
    pdf.add_font("DejaVu", "", FONT_PATH)
    pdf.set_font("DejaVu", size=16)
    pdf.cell(0, 12, "Карточки для повторения", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(4)

    for card in cards:
        pdf.set_font("DejaVu", size=13)
        pdf.multi_cell(0, 8, card.get("title", ""), new_x="LMARGIN", new_y="NEXT")

        pdf.set_font("DejaVu", size=10)
        tags = " · ".join(card.get("tags", []))
        if tags:
            pdf.multi_cell(0, 6, tags, new_x="LMARGIN", new_y="NEXT")

        pdf.set_font("DejaVu", size=11)
        pdf.multi_cell(0, 7, "Вопрос: " + card.get("question", ""), new_x="LMARGIN", new_y="NEXT")
        pdf.multi_cell(0, 7, "Ответ: " + card.get("answer", ""), new_x="LMARGIN", new_y="NEXT")
        pdf.ln(6)

    buffer = io.BytesIO(pdf.output())
    buffer.seek(0)
    return buffer


@app.get("/")
async def root():
    return FileResponse(os.path.join("static", "index.html"))


app.mount("/static", StaticFiles(directory="static"), name="static")
